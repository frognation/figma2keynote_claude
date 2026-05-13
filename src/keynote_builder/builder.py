"""
Keynote .key file builder.

Strategy:
  1. Unpack a template .key file using keynote-parser
  2. Modify the YAML/protobuf to add slides, text, images, videos
  3. Repack into a valid .key file

The .key format:
  - ZIP archive containing:
    - Index/Slide-*.iwa  (Snappy-compressed Protobuf)
    - Data/              (media files: images, videos)
    - Metadata/          (document-level metadata)
    - preview.jpg        (thumbnail)
"""

import os
import json
import shutil
import zipfile
import struct
from pathlib import Path
from typing import Optional
from copy import deepcopy

# keynote-parser imports
try:
    from keynote_parser import file_utils
    HAS_KEYNOTE_PARSER = True
except ImportError:
    HAS_KEYNOTE_PARSER = False


class KeynoteBuilder:
    """
    Builds a .key file from a manifest + assets.

    Two modes:
      1. Template mode: modify an existing .key template
      2. Raw mode: build .key from scratch using protobuf (advanced)
    """

    def __init__(self, template_path: Optional[Path] = None):
        self.template_path = Path(template_path) if template_path else None
        self.work_dir = None

    def build_from_manifest(
        self,
        manifest: dict,
        assets_dir: Path,
        output_path: Path,
    ) -> Path:
        """
        Main entry: manifest + assets → .key file.

        Falls back to PPTX-based approach if no template available.
        """
        if self.template_path and self.template_path.exists():
            return self._build_from_template(manifest, assets_dir, output_path)
        else:
            return self._build_via_pptx(manifest, assets_dir, output_path)

    # ── PPTX-based approach (more reliable, works without template) ──

    def _build_via_pptx(
        self, manifest: dict, assets_dir: Path, output_path: Path
    ) -> Path:
        """
        Build a .pptx file that Keynote can open natively.

        This is the pragmatic MVP approach:
        - python-pptx creates fully editable .pptx
        - Keynote opens .pptx with editable text
        - Videos embedded as media clips
        - Then optionally convert to .key via AppleScript
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ImportError(
                "python-pptx required for PPTX mode. "
                "Install with: pip install python-pptx"
            )

        # Create presentation with correct dimensions
        canvas = manifest.get("canvas", {})
        width = canvas.get("width", 1920)
        height = canvas.get("height", 1080)

        prs = Presentation()
        prs.slide_width = Emu(int(width * 914400 / 96))   # px to EMU
        prs.slide_height = Emu(int(height * 914400 / 96))

        # Blank layout
        blank_layout = prs.slide_layouts[6]  # Typically blank

        for slide_data in manifest.get("slides", []):
            slide = prs.slides.add_slide(blank_layout)

            # Set background
            bg = slide_data.get("background", {})
            if bg.get("type") == "solid":
                color = bg.get("color", {})
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(
                    int(color.get("r", 1) * 255),
                    int(color.get("g", 1) * 255),
                    int(color.get("b", 1) * 255),
                )

            # Process elements
            self._add_elements_to_pptx_slide(
                slide, slide_data.get("elements", []),
                assets_dir, manifest
            )

        # Save
        pptx_path = output_path.with_suffix(".pptx")
        prs.save(str(pptx_path))

        # Attempt auto-convert to .key via AppleScript
        key_path = self._convert_pptx_to_key(pptx_path, output_path)

        return key_path or pptx_path

    def _add_elements_to_pptx_slide(
        self, slide, elements: list, assets_dir: Path, manifest: dict
    ):
        """Add elements from manifest to a pptx slide."""
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        scale = 914400 / 96  # px → EMU conversion

        for el in elements:
            el_type = el.get("type", "")
            x = int(el.get("x", 0) * scale)
            y = int(el.get("y", 0) * scale)
            w = int(el.get("width", 100) * scale)
            h = int(el.get("height", 100) * scale)

            if el_type == "TEXT" and el.get("text_data"):
                self._add_text_box(slide, el, x, y, w, h)

            elif el.get("media_data"):
                media = el["media_data"]
                if media["type"] == "image" and media.get("local_path"):
                    img_path = assets_dir.parent / media["local_path"]
                    if img_path.exists():
                        slide.shapes.add_picture(
                            str(img_path), Emu(x), Emu(y), Emu(w), Emu(h)
                        )

                elif media["type"] == "video" and media.get("local_path"):
                    vid_path = assets_dir.parent / media["local_path"]
                    if vid_path.exists():
                        # Get poster image if available
                        poster = None
                        slide_idx = el.get("_slide_index", 0)
                        poster_candidates = list(
                            assets_dir.glob(f"slide_{slide_idx:03d}_poster.*")
                        )
                        if poster_candidates:
                            poster = str(poster_candidates[0])

                        self._add_video(
                            slide, str(vid_path), poster,
                            Emu(x), Emu(y), Emu(w), Emu(h)
                        )

            elif el_type in ("RECTANGLE", "ELLIPSE") and el.get("shape_data"):
                self._add_shape(slide, el, x, y, w, h)

            # Recurse into children
            if el.get("children"):
                self._add_elements_to_pptx_slide(
                    slide, el["children"], assets_dir, manifest
                )

    def _add_text_box(self, slide, el: dict, x: int, y: int, w: int, h: int):
        """Add an editable text box to the slide."""
        from pptx.util import Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        text_data = el["text_data"]
        txBox = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
        tf = txBox.text_frame
        tf.word_wrap = True

        # Base style
        style = text_data.get("style", {})

        # Alignment mapping
        align_map = {
            "LEFT": PP_ALIGN.LEFT,
            "CENTER": PP_ALIGN.CENTER,
            "RIGHT": PP_ALIGN.RIGHT,
            "JUSTIFIED": PP_ALIGN.JUSTIFY,
        }

        # Handle character-level styling
        characters = text_data.get("characters", "")
        char_overrides = text_data.get("characterStyleOverrides", [])
        style_table = text_data.get("styleOverrideTable", {})

        # Split text into paragraphs
        paragraphs = characters.split("\n")
        char_pos = 0

        for p_idx, para_text in enumerate(paragraphs):
            if p_idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.alignment = align_map.get(
                style.get("textAlignHorizontal", "LEFT"), PP_ALIGN.LEFT
            )

            if not para_text:
                char_pos += 1  # newline
                continue

            # Build runs based on character style overrides
            runs = self._build_styled_runs(
                para_text, char_pos, char_overrides, style_table, style
            )

            for run_text, run_style in runs:
                run = p.add_run()
                run.text = run_text

                # Apply font properties
                font = run.font
                font.name = run_style.get("fontFamily", style.get("fontFamily", "Helvetica"))
                raw_size = run_style.get("fontSize", style.get("fontSize", 16))
                font.size = Pt(max(raw_size, 1))  # min 1pt (pptx requires >= ~0.8pt)
                font.bold = run_style.get("fontWeight", style.get("fontWeight", 400)) >= 700
                font.italic = run_style.get("italic", style.get("italic", False))

                # Text color
                fills = run_style.get("fills", text_data.get("fills", []))
                if fills and fills[0].get("type") == "SOLID":
                    c = fills[0]["color"]
                    font.color.rgb = RGBColor(
                        int(c.get("r", 0) * 255),
                        int(c.get("g", 0) * 255),
                        int(c.get("b", 0) * 255),
                    )

                # Underline / strikethrough
                decoration = run_style.get(
                    "textDecoration", style.get("textDecoration", "NONE")
                )
                if decoration == "UNDERLINE":
                    font.underline = True
                elif decoration == "STRIKETHROUGH":
                    font.strikethrough = True

            char_pos += len(para_text) + 1  # +1 for newline

    def _build_styled_runs(
        self, text: str, start_pos: int,
        char_overrides: list, style_table: dict, base_style: dict
    ) -> list[tuple[str, dict]]:
        """
        Split text into runs based on character style overrides.
        Returns [(text, merged_style), ...]
        """
        if not char_overrides:
            return [(text, base_style)]

        runs = []
        current_style_id = None
        current_text = ""

        for i, ch in enumerate(text):
            global_pos = start_pos + i
            if global_pos < len(char_overrides):
                style_id = str(char_overrides[global_pos])
            else:
                style_id = "0"  # base style

            if style_id != current_style_id:
                if current_text:
                    merged = {**base_style}
                    if current_style_id and current_style_id in style_table:
                        merged.update(style_table[current_style_id])
                    runs.append((current_text, merged))
                current_text = ch
                current_style_id = style_id
            else:
                current_text += ch

        if current_text:
            merged = {**base_style}
            if current_style_id and current_style_id in style_table:
                merged.update(style_table[current_style_id])
            runs.append((current_text, merged))

        return runs if runs else [(text, base_style)]

    def _add_shape(self, slide, el: dict, x: int, y: int, w: int, h: int):
        """Add a basic shape."""
        from pptx.util import Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        shape_type = MSO_SHAPE.RECTANGLE
        if el.get("type") == "ELLIPSE":
            shape_type = MSO_SHAPE.OVAL

        shape = slide.shapes.add_shape(
            shape_type, Emu(x), Emu(y), Emu(w), Emu(h)
        )

        shape_data = el.get("shape_data", {})
        fills = shape_data.get("fills", [])

        if fills and fills[0].get("type") == "SOLID":
            c = fills[0]["color"]
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(
                int(c.get("r", 0) * 255),
                int(c.get("g", 0) * 255),
                int(c.get("b", 0) * 255),
            )

        # Corner radius
        cr = shape_data.get("cornerRadius", 0)
        if cr > 0:
            # python-pptx doesn't directly support corner radius on rectangles
            # but we can use ROUNDED_RECTANGLE
            pass

        # Stroke
        sw = shape_data.get("strokeWeight", 0)
        strokes = shape_data.get("strokes", [])
        if sw > 0 and strokes:
            from pptx.util import Pt
            shape.line.width = Pt(sw)
            if strokes[0].get("type") == "SOLID":
                sc = strokes[0]["color"]
                shape.line.color.rgb = RGBColor(
                    int(sc.get("r", 0) * 255),
                    int(sc.get("g", 0) * 255),
                    int(sc.get("b", 0) * 255),
                )
        else:
            shape.line.fill.background()  # No border

    def _add_video(
        self, slide, video_path: str, poster_path: Optional[str],
        left, top, width, height
    ):
        """
        Embed a video into the slide.

        python-pptx supports video embedding via add_movie.
        """
        from pptx.util import Emu

        if poster_path and os.path.exists(poster_path):
            slide.shapes.add_movie(
                video_path, left, top, width, height,
                poster_frame_image=poster_path,
            )
        else:
            # Create a placeholder poster (black frame)
            from PIL import Image
            poster_tmp = Path(video_path).parent / "_poster_tmp.png"
            img = Image.new("RGB", (1920, 1080), (0, 0, 0))
            img.save(str(poster_tmp))
            slide.shapes.add_movie(
                video_path, left, top, width, height,
                poster_frame_image=str(poster_tmp),
            )

    # ── PPTX → .key conversion via AppleScript ──────────────────

    # ── Keynote app targeting ────────────────────────────────────
    # On systems with multiple Keynote versions installed, we pin to the
    # specific newer version by bundle ID. The two known bundle IDs:
    #   com.apple.iWork.Keynote → Keynote v14.x (older, original name)
    #   com.apple.Keynote       → Keynote v15.x ("Keynote Creator Studio")
    # We prefer the newer version.
    KEYNOTE_BUNDLE_ID_NEW = "com.apple.Keynote"           # v15+
    KEYNOTE_BUNDLE_ID_OLD = "com.apple.iWork.Keynote"     # v14.x

    @classmethod
    def detect_best_keynote(cls) -> str:
        """Find the best (newest) installed Keynote and return its bundle ID."""
        import subprocess
        # Try newer first, fall back to older
        for bundle_id in (cls.KEYNOTE_BUNDLE_ID_NEW, cls.KEYNOTE_BUNDLE_ID_OLD):
            try:
                result = subprocess.run(
                    ["osascript", "-e",
                     f'tell application id "{bundle_id}" to version'],
                    capture_output=True, text=True, timeout=5,
                )
                if result.returncode == 0 and result.stdout.strip():
                    return bundle_id
            except (subprocess.TimeoutExpired, FileNotFoundError):
                continue
        return cls.KEYNOTE_BUNDLE_ID_NEW  # default

    def _convert_pptx_to_key(self, pptx_path: Path, key_path: Path) -> Optional[Path]:
        """
        Convert .pptx to .key using Keynote via AppleScript.
        Pins to the specific Keynote bundle ID to avoid version conflicts.
        """
        import subprocess

        bundle_id = self.detect_best_keynote()

        script = f'''
        tell application id "{bundle_id}"
            activate
            delay 1
            open POSIX file "{pptx_path}"
            delay 5
            set theDoc to front document
            save theDoc in POSIX file "{key_path}"
            delay 2
            close theDoc saving no
        end tell
        '''

        try:
            result = subprocess.run(
                ["osascript", "-e", script],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0 and key_path.exists():
                return key_path
            elif result.stderr:
                print(f"  [AppleScript] {result.stderr.strip()}")
        except (subprocess.TimeoutExpired, FileNotFoundError) as e:
            print(f"  [AppleScript] {e}")

        return None

    # ── Template-based approach (for native .key) ────────────────

    def _build_from_template(
        self, manifest: dict, assets_dir: Path, output_path: Path
    ) -> Path:
        """
        Build .key by modifying a template.
        Uses keynote-parser to unpack/repack.
        """
        if not HAS_KEYNOTE_PARSER:
            raise ImportError("keynote-parser required for template mode")

        # 1. Copy template
        work_key = output_path.with_suffix(".key.tmp")
        shutil.copy2(self.template_path, work_key)

        # 2. Unpack
        unpack_dir = output_path.with_suffix(".key.unpacked")
        if unpack_dir.exists():
            shutil.rmtree(unpack_dir)
        file_utils.unpack(str(work_key), str(unpack_dir))

        # 3. Copy media assets into Data/
        data_dir = unpack_dir / "Data"
        data_dir.mkdir(exist_ok=True)

        asset_mapping = {}  # local_path → Data/ filename
        for asset_file in assets_dir.iterdir():
            if asset_file.suffix.lower() in (
                ".png", ".jpg", ".jpeg", ".gif", ".mp4", ".mov", ".m4v", ".webm"
            ):
                dest = data_dir / asset_file.name
                shutil.copy2(asset_file, dest)
                asset_mapping[str(asset_file)] = asset_file.name

        # 4. Repack
        file_utils.pack(str(unpack_dir), str(output_path))

        # 5. Cleanup
        if work_key.exists():
            work_key.unlink()
        if unpack_dir.exists():
            shutil.rmtree(unpack_dir)

        return output_path
